from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


NAVY = RGBColor(8, 35, 74)
TEAL = RGBColor(0, 132, 137)
AMBER = RGBColor(240, 164, 35)
SLATE = RGBColor(51, 65, 85)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(107, 114, 128)

FONT = "Trebuchet MS"


def add_notes(slide, text: str) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = text


def set_text_style(paragraph, size: int, color: RGBColor = SLATE, bold: bool = False) -> None:
    paragraph.font.name = FONT
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color


def add_top_band(slide, title: str, subtitle: str = "") -> None:
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.92),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.16), Inches(8.8), Inches(0.45))
    title_tf = title_box.text_frame
    title_tf.clear()
    title_p = title_tf.paragraphs[0]
    title_p.text = title
    set_text_style(title_p, 24, WHITE, True)

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.53), Inches(8.8), Inches(0.25))
        sub_tf = sub_box.text_frame
        sub_tf.clear()
        sub_p = sub_tf.paragraphs[0]
        sub_p.text = subtitle
        set_text_style(sub_p, 12, RGBColor(219, 234, 254), False)

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(11.9),
        Inches(0),
        Inches(1.433),
        Inches(0.92),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()


def add_footer(slide, text: str = "SBCbank | Consultant Proposal | Confidential") -> None:
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.4),
        Inches(7.05),
        Inches(12.533),
        Inches(0.01),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(203, 213, 225)
    line.line.fill.background()

    footer = slide.shapes.add_textbox(Inches(0.45), Inches(7.08), Inches(12.4), Inches(0.22))
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    set_text_style(p, 10, MUTED, False)


def add_bullets(slide, left: float, top: float, width: float, height: float, bullets: list[str], size: int = 20) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    for index, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = f"- {bullet}"
        set_text_style(p, size, SLATE, False)
        p.space_after = Pt(10)


def build_deck(repo_root: Path) -> Path:
    output_dir = repo_root / "artifacts"
    output_dir.mkdir(parents=True, exist_ok=True)

    arch_path = output_dir / "sbcbank-architecture-full.png"
    iam_path = output_dir / "iam-diagram-organized.png"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(6.5), Inches(13.333), Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8), Inches(3.8), Inches(0.6))
    tag.fill.solid()
    tag.fill.fore_color.rgb = AMBER
    tag.line.fill.background()
    tag_tf = tag.text_frame
    tag_tf.clear()
    tag_p = tag_tf.paragraphs[0]
    tag_p.text = "CONSULTANT PITCH"
    set_text_style(tag_p, 14, NAVY, True)
    tag_p.alignment = PP_ALIGN.CENTER

    title = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.8), Inches(2.0))
    title_tf = title.text_frame
    title_tf.clear()
    p = title_tf.paragraphs[0]
    p.text = "SBCbank Cloud-Native Digital Bank"
    set_text_style(p, 44, WHITE, True)
    p = title_tf.add_paragraph()
    p.text = "Board and Owner Investment Proposal"
    set_text_style(p, 24, RGBColor(191, 219, 254), False)

    meta = slide.shapes.add_textbox(Inches(0.8), Inches(5.95), Inches(11.8), Inches(0.45))
    meta_tf = meta.text_frame
    meta_tf.clear()
    m = meta_tf.paragraphs[0]
    m.text = "April 2026 | Singapore (MAS/PDPA-Aligned)"
    set_text_style(m, 16, WHITE, False)

    add_notes(
        slide,
        "Open with business outcomes: profitable growth, operational resilience, and compliance confidence. "
        "Set expectation that this is an execution-ready plan, not only a technical concept.",
    )

    # Slide 2: Owner priorities
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_band(slide, "What The Bank Owner Needs To Hear", "Outcome-first framing")
    add_bullets(
        slide,
        0.8,
        1.35,
        12.0,
        3.8,
        [
            "How fast we can acquire and retain profitable customers.",
            "How the platform lowers cost-to-serve while scaling safely.",
            "How we remain MAS/PDPA compliant by design, not retrofitted.",
            "How we manage fraud, outages, and reputational risk.",
            "How this becomes a durable competitive advantage in 12-24 months.",
        ],
        size=24,
    )

    kpi_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.25), Inches(12.0), Inches(1.55))
    kpi_box.fill.solid()
    kpi_box.fill.fore_color.rgb = RGBColor(239, 246, 255)
    kpi_box.line.color.rgb = RGBColor(125, 211, 252)

    kpi_tf = kpi_box.text_frame
    kpi_tf.clear()
    p = kpi_tf.paragraphs[0]
    p.text = "Board-level headline metrics"
    set_text_style(p, 16, NAVY, True)
    p = kpi_tf.add_paragraph()
    p.text = "Time to MVP, cost per transaction, fraud loss rate, uptime SLA, and compliance audit pass rate."
    set_text_style(p, 14, SLATE, False)

    add_footer(slide)
    add_notes(slide, "Frame the conversation around certainty: returns, risk controls, and execution credibility.")

    # Slide 3: Solution summary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_band(slide, "Solution Summary", "Cloud-native core designed for growth and control")

    left_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.8), Inches(1.35), Inches(6.0), Inches(5.45))
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = RGBColor(248, 250, 252)
    left_panel.line.color.rgb = RGBColor(203, 213, 225)

    right_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(6.95), Inches(1.35), Inches(5.55), Inches(5.45))
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = RGBColor(240, 253, 250)
    right_panel.line.color.rgb = RGBColor(94, 234, 212)

    add_bullets(
        slide,
        1.1,
        1.7,
        5.5,
        4.9,
        [
            "Microservices for account, payment, ledger, statement, notification, and orchestration.",
            "Event-driven workflow using Step Functions, EventBridge, and SQS for operational resilience.",
            "PostgreSQL + DynamoDB ledger patterns for transactional integrity and traceability.",
            "Cloud-native controls: encryption, least-privilege IAM, centralized logging.",
        ],
        size=16,
    )

    right_text = slide.shapes.add_textbox(Inches(7.25), Inches(1.7), Inches(4.9), Inches(4.9))
    tf = right_text.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Business impact"
    set_text_style(p, 20, NAVY, True)
    for line in [
        "- Faster product launch cycles",
        "- Lower change failure risk",
        "- Clear operational accountability",
        "- Scalable cost profile",
        "- Better regulator and auditor readiness",
    ]:
        p = tf.add_paragraph()
        p.text = line
        set_text_style(p, 16, SLATE, False)
        p.space_after = Pt(8)

    add_footer(slide)
    add_notes(slide, "Position architecture as a means to business outcomes, not as an engineering vanity project.")

    # Slide 4: Architecture diagram
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_band(slide, "Reference Architecture", "Derived from existing sbcbank-architecture-full.drawio.xml")

    if arch_path.exists():
        slide.shapes.add_picture(str(arch_path), Inches(0.75), Inches(1.25), width=Inches(12.0))
    else:
        placeholder = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.75), Inches(1.25), Inches(12.0), Inches(4.9))
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(254, 242, 242)
        placeholder.line.color.rgb = RGBColor(252, 165, 165)
        ph_tf = placeholder.text_frame
        ph_tf.text = "Architecture diagram image not found in artifacts/."

    add_bullets(
        slide,
        0.9,
        6.2,
        11.8,
        0.7,
        ["Flow: Channel layer -> API and services -> orchestration and events -> data and observability controls."],
        size=13,
    )

    add_footer(slide)
    add_notes(slide, "Walk left to right. Emphasize isolation, fault containment, and scale-out characteristics.")

    # Slide 5: IAM and compliance diagram
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_band(slide, "Security and Compliance Control Plane", "Derived from existing iam-diagram-organized.drawio.xml")

    if iam_path.exists():
        slide.shapes.add_picture(str(iam_path), Inches(0.75), Inches(1.25), width=Inches(12.0))
    else:
        placeholder = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.75), Inches(1.25), Inches(12.0), Inches(4.9))
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(254, 242, 242)
        placeholder.line.color.rgb = RGBColor(252, 165, 165)
        ph_tf = placeholder.text_frame
        ph_tf.text = "IAM diagram image not found in artifacts/."

    add_bullets(
        slide,
        0.9,
        6.15,
        12.0,
        0.8,
        ["Least privilege IAM, KMS key separation, and centralized logging align with MAS TRM and PDPA obligations."],
        size=13,
    )

    add_footer(slide)
    add_notes(slide, "Highlight regulator confidence: traceability, role separation, and encryption boundaries.")

    # Slide 6: Customer and market value
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_band(slide, "Customer and Market Proposition", "How this wins in Singapore")

    cards = [
        ("Retail Customers", "Instant onboarding, reliable transfers, transparent statements, fast dispute support."),
        ("Merchants", "Faster settlement visibility, fewer payment failures, cleaner reconciliation data."),
        ("Bank Leadership", "Launch products faster while reducing operational and compliance drag."),
    ]

    x = 0.8
    for title, body in cards:
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.6), Inches(4.1), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(248, 250, 252)
        card.line.color.rgb = RGBColor(203, 213, 225)

        tf = card.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        set_text_style(p, 18, NAVY, True)
        p = tf.add_paragraph()
        p.text = body
        set_text_style(p, 14, SLATE, False)
        p.space_before = Pt(8)
        x += 4.25

    add_bullets(
        slide,
        0.8,
        6.25,
        12.2,
        0.5,
        ["Differentiator: resilient architecture paired with compliance automation, not just feature parity."],
        size=13,
    )

    add_footer(slide)
    add_notes(slide, "Anchor on differentiated trust and reliability, not only UI or pricing.")

    # Slide 7: Financial model with chart
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_band(slide, "Illustrative 3-Year Value Model", "For board decision support; refine with finance team")

    assumptions = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.8), Inches(1.35), Inches(4.2), Inches(5.3))
    assumptions.fill.solid()
    assumptions.fill.fore_color.rgb = RGBColor(255, 251, 235)
    assumptions.line.color.rgb = RGBColor(253, 224, 71)

    tf = assumptions.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Model assumptions"
    set_text_style(p, 17, NAVY, True)
    for line in [
        "- MVP launch in 4-6 months",
        "- Progressive onboarding ramp",
        "- Fraud loss rate reduction via workflow controls",
        "- Cloud spend governance from month 1",
        "- Benefits include cost and revenue effects",
    ]:
        p = tf.add_paragraph()
        p.text = line
        set_text_style(p, 13, SLATE, False)
        p.space_after = Pt(6)

    chart_data = CategoryChartData()
    chart_data.categories = ["Year 1", "Year 2", "Year 3"]
    chart_data.add_series("Investment (M SGD)", (1.2, 0.9, 1.0))
    chart_data.add_series("Estimated Benefits (M SGD)", (0.6, 1.6, 2.8))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(5.3),
        Inches(1.55),
        Inches(7.2),
        Inches(4.8),
        chart_data,
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.category_axis.has_major_gridlines = False
    chart.value_axis.has_major_gridlines = True

    add_bullets(
        slide,
        5.3,
        6.35,
        7.2,
        0.5,
        ["Indicative break-even expected during Year 2 with upside from new product velocity."],
        size=12,
    )

    add_footer(slide)
    add_notes(slide, "State clearly these numbers are illustrative and should be pressure-tested with internal finance assumptions.")

    # Slide 8: Delivery roadmap
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_band(slide, "12-Month Delivery Roadmap", "Milestone-based execution with clear gates")

    phases = [
        ("Phase 1", "Months 0-2", "Foundation\nControls\nCore services"),
        ("Phase 2", "Months 3-5", "MVP launch\nPayments\nStatements"),
        ("Phase 3", "Months 6-9", "Fraud workflow\nOperational tooling\nDashboards"),
        ("Phase 4", "Months 10-12", "Scale and optimize\nCost tuning\nAudit hardening"),
    ]

    x = 0.8
    for idx, (phase, months, detail) in enumerate(phases):
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.1), Inches(2.85), Inches(3.0))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(239, 246, 255) if idx % 2 == 0 else RGBColor(236, 253, 245)
        box.line.color.rgb = RGBColor(125, 211, 252) if idx % 2 == 0 else RGBColor(110, 231, 183)

        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = phase
        set_text_style(p, 16, NAVY, True)
        p.alignment = PP_ALIGN.CENTER
        p = tf.add_paragraph()
        p.text = months
        set_text_style(p, 13, TEAL, True)
        p.alignment = PP_ALIGN.CENTER
        p = tf.add_paragraph()
        p.text = detail
        set_text_style(p, 12, SLATE, False)
        p.alignment = PP_ALIGN.CENTER

        if idx < len(phases) - 1:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x + 2.95), Inches(3.2), Inches(0.45), Inches(0.7))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(148, 163, 184)
            arrow.line.fill.background()
        x += 3.1

    add_bullets(
        slide,
        0.9,
        5.5,
        12.1,
        1.2,
        [
            "Decision gates at end of each phase: risk, budget, compliance readiness, and customer outcomes.",
            "No production expansion without passing security and operational readiness criteria.",
        ],
        size=14,
    )

    add_footer(slide)
    add_notes(slide, "This timeline gives the board control over investment pacing and risk exposure.")

    # Slide 9: Risks and mitigations
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_band(slide, "Top Risks and Mitigations", "What could fail and how we prevent it")

    table_shape = slide.shapes.add_table(6, 2, Inches(0.8), Inches(1.45), Inches(12.0), Inches(5.6)).table
    table_shape.columns[0].width = Inches(4.0)
    table_shape.columns[1].width = Inches(8.0)

    headers = ["Risk", "Mitigation"]
    for col, head in enumerate(headers):
        cell = table_shape.cell(0, col)
        cell.text = head
        for p in cell.text_frame.paragraphs:
            set_text_style(p, 14, WHITE, True)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY

    rows = [
        ("Regulatory breach", "Policy-as-code controls, encrypted data planes, and auditable log retention."),
        ("Fraud spike", "Step Functions fraud branch, manual review queues, and incident runbooks."),
        ("Cost overrun", "Tagging, budgets, FinOps reviews, and environment-level guardrails."),
        ("Delivery slippage", "Milestone gating, scope discipline, and weekly executive steering."),
        ("Service outage", "Multi-AZ targets, failover tests, and controlled rollback procedures."),
    ]

    for row_idx, (risk, mitigation) in enumerate(rows, start=1):
        cell_risk = table_shape.cell(row_idx, 0)
        cell_risk.text = risk
        for p in cell_risk.text_frame.paragraphs:
            set_text_style(p, 13, NAVY, True)
        cell_risk.fill.solid()
        cell_risk.fill.fore_color.rgb = RGBColor(241, 245, 249)

        cell_mit = table_shape.cell(row_idx, 1)
        cell_mit.text = mitigation
        for p in cell_mit.text_frame.paragraphs:
            set_text_style(p, 12, SLATE, False)
        cell_mit.fill.solid()
        cell_mit.fill.fore_color.rgb = RGBColor(255, 255, 255)

    add_footer(slide)
    add_notes(slide, "Show that risk is expected, measured, and actively controlled through architecture and governance.")

    # Slide 10: Decision and ask
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_band(slide, "Decision Required", "Approve launch of Phase 1 with governance guardrails")

    ask_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(12.0), Inches(2.3))
    ask_box.fill.solid()
    ask_box.fill.fore_color.rgb = RGBColor(236, 253, 245)
    ask_box.line.color.rgb = RGBColor(52, 211, 153)

    tf = ask_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Requested approval"
    set_text_style(p, 20, NAVY, True)
    p = tf.add_paragraph()
    p.text = "Authorize Phase 1 funding, assign executive sponsor, and establish monthly board checkpoints."
    set_text_style(p, 16, SLATE, False)

    add_bullets(
        slide,
        0.9,
        4.15,
        12.0,
        2.3,
        [
            "Budget envelope: confirm with finance model after assumption workshop.",
            "Governance: security, risk, and product leadership included in gate reviews.",
            "Immediate next action: run a 2-week discovery sprint to lock execution baseline.",
        ],
        size=18,
    )

    close = slide.shapes.add_textbox(Inches(0.8), Inches(6.45), Inches(12.0), Inches(0.45))
    close_tf = close.text_frame
    close_tf.clear()
    p = close_tf.paragraphs[0]
    p.text = "This proposal is designed to deliver growth with control, not growth with hidden risk."
    set_text_style(p, 14, NAVY, True)
    p.alignment = PP_ALIGN.CENTER

    add_footer(slide)
    add_notes(slide, "Close with confidence and a clear ask. Leave no ambiguity on next decision.")

    output_path = output_dir / "SBCbank_BankOwner_Pitch_Deck_Apr2026.pptx"
    prs.save(str(output_path))
    return output_path


def main() -> None:
    repo_root = Path(__file__).resolve().parents[1]
    output_path = build_deck(repo_root)
    print(f"Created deck: {output_path}")


if __name__ == "__main__":
    main()
